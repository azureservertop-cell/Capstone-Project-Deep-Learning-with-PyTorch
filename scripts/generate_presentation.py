from pathlib import Path
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = Path('presentation/EuroSAT_Capstone_Presentation.pptx')
OUT.parent.mkdir(parents=True, exist_ok=True)
ASSET_DIR = Path('presentation/generated_assets')
ASSET_DIR.mkdir(parents=True, exist_ok=True)

BLUE = RGBColor(38, 101, 214)
DARK = RGBColor(15, 23, 42)
GRAY = RGBColor(75, 85, 99)
LIGHT = RGBColor(239, 246, 255)
GREEN = RGBColor(220, 252, 231)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 250, 252)
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.09), prs.slide_height)
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.color.rgb = BLUE


def tx(slide, text, x, y, w, h, size=18, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = 'Aptos'
    return box


def title(slide, t, sub=None):
    tx(slide, t, 0.35, 0.35, 7.7, 0.45, 24, True)
    if sub:
        tx(slide, sub, 0.35, 0.84, 8.6, 0.3, 9, False, GRAY)


def metric(slide, x, y, w, h, top, bottom, fill=LIGHT):
    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = fill; rect.line.color.rgb = RGBColor(219, 234, 254)
    tx(slide, top, x, y+0.2, w, 0.28, 8, False, GRAY, PP_ALIGN.CENTER)
    tx(slide, bottom, x, y+0.52, w, 0.34, 15, True, DARK, PP_ALIGN.CENTER)


def bullet(slide, items, x, y, w, h, size=12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = 'Aptos'
        p.font.color.rgb = DARK
        p.text = u'• ' + it
    return box


def footer(slide, n):
    tx(slide, f'EuroSAT Capstone | {n}', 11.0, 7.14, 1.8, 0.18, 6, False, GRAY)


def add_chart_pipeline(path):
    labels=['Baseline','Incremental Final']
    acc=[85.60,92.07]
    f1=[85.49,92.09]
    x=range(len(labels))
    fig, ax=plt.subplots(figsize=(5.2,3.2))
    ax.bar([i-0.18 for i in x], acc, width=0.35, label='Accuracy')
    ax.bar([i+0.18 for i in x], f1, width=0.35, label='F1-score')
    ax.set_ylim(0,100); ax.set_ylabel('Percent'); ax.set_title('Baseline vs Incremental Training Setup', fontsize=10)
    ax.set_xticks(list(x)); ax.set_xticklabels(labels, fontsize=8)
    ax.legend(fontsize=8); ax.grid(axis='y', alpha=0.25)
    fig.tight_layout(); fig.savefig(path, dpi=200); plt.close(fig)


def add_chart_tuning(path):
    labels=['Exp 1\nLR .001\nB32 D.3','Exp 2\nLR .0005\nB32 D.3','Exp 3\nLR .001\nB64 D.5']
    vals=[0.4675,0.4457,0.6175]
    fig, ax=plt.subplots(figsize=(5.1,3.2))
    ax.bar(labels, vals)
    ax.set_ylabel('Validation Loss'); ax.set_title('6-Epoch Hyperparameter Tuning', fontsize=10)
    ax.grid(axis='y', alpha=0.25)
    for i,v in enumerate(vals): ax.text(i, v+0.012, f'{v:.4f}', ha='center', fontsize=7)
    fig.tight_layout(); fig.savefig(path, dpi=200); plt.close(fig)


def add_chart_compare(path):
    labels=['Final Custom CNN','Frozen ResNet18']
    acc=[92.07,82.44]; f1=[92.09,82.32]
    x=range(len(labels))
    fig, ax=plt.subplots(figsize=(5.2,3.2))
    ax.bar([i-0.18 for i in x], acc, width=0.35, label='Accuracy')
    ax.bar([i+0.18 for i in x], f1, width=0.35, label='F1-score')
    ax.set_ylim(0,100); ax.set_ylabel('Percent'); ax.set_xticks(list(x)); ax.set_xticklabels(labels, fontsize=8)
    ax.set_title('Final Model Comparison', fontsize=10); ax.legend(fontsize=8); ax.grid(axis='y', alpha=0.25)
    fig.tight_layout(); fig.savefig(path, dpi=200); plt.close(fig)


def add_chart_ablation(path):
    labels=['Final baseline','+ ColorJitter','+ VerticalFlip','+ Rotation pipeline']
    vals=[92.07,91.48,89.83,67.60]
    fig, ax=plt.subplots(figsize=(5.2,3.2))
    ax.bar(labels, vals)
    ax.set_ylim(60,100); ax.set_ylabel('Test Accuracy (%)'); ax.set_title('Extra Augmentation Did Not Improve Accuracy', fontsize=10)
    ax.tick_params(axis='x', labelrotation=25, labelsize=7); ax.grid(axis='y', alpha=0.25)
    fig.tight_layout(); fig.savefig(path, dpi=200); plt.close(fig)

pipeline = ASSET_DIR/'pipeline.png'; add_chart_pipeline(pipeline)
tuning = ASSET_DIR/'tuning.png'; add_chart_tuning(tuning)
compare = ASSET_DIR/'compare.png'; add_chart_compare(compare)
ablation = ASSET_DIR/'ablation.png'; add_chart_ablation(ablation)

# 1
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'EuroSAT Satellite Image Classification','Deep Learning Capstone Project using PyTorch')
metric(s,0.55,2.35,1.45,0.8,'Final CNN Accuracy','92.07%')
metric(s,2.3,2.35,1.45,0.8,'Weighted F1-score','92.09%',GREEN)
metric(s,4.05,2.35,1.0,0.8,'Classes','10')
tx(s,'Jianye Chen | Fanshawe College',0.55,6.55,2.6,0.25,7,False,GRAY); footer(s,1)

# 2
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Task and Dataset','Classify EuroSAT satellite images into 10 land-use categories.')
metric(s,0.55,1.3,1.4,0.75,'Total images','27,000'); metric(s,2.15,1.3,1.6,0.75,'Train / Val / Test','70 / 15 / 15',GREEN)
bullet(s,['RGB remote-sensing images','10 balanced land-use classes','Stratified split to prevent class imbalance artifacts','Test set held out until final evaluation'],0.55,2.4,4.4,1.5,11)
if Path('images/class_distribution.png').exists(): s.shapes.add_picture('images/class_distribution.png', Inches(6.35), Inches(0.9), width=Inches(5.3))
if Path('images/augmented_samples.png').exists(): s.shapes.add_picture('images/augmented_samples.png', Inches(6.35), Inches(3.7), width=Inches(4.5))
footer(s,2)

# 3 workflow
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Project Workflow','The project follows a complete deep-learning experiment pipeline.')
steps=[('Dataset','EuroSAT'),('Preprocess','Resize + normalize + augment'),('CNN','3-layer CNN'),('Tuning','Manual 3-setting comparison'),('Evaluate','Metrics + confusion matrix'),('Ablation','Test extra augmentation')]
x=0.8
for a,b in steps:
    rect=s.shapes.add_shape(1, Inches(x), Inches(2.1), Inches(1.5), Inches(0.8)); rect.fill.solid(); rect.fill.fore_color.rgb=LIGHT; rect.line.color.rgb=RGBColor(191,219,254)
    tx(s,a,x,2.28,1.5,0.2,9,True,DARK,PP_ALIGN.CENTER); tx(s,b,x,2.52,1.5,0.2,6,False,GRAY,PP_ALIGN.CENTER)
    x+=1.95
bullet(s,['The final model was selected using validation loss.','The final score was measured on a held-out test set.','Ablation results justify why extra preprocessing was not adopted.'],0.8,4.1,8.5,0.9,11)
footer(s,3)

# 4 preprocessing
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Final Preprocessing Choice','Simple augmentation performed best and kept the pipeline focused.')
bullet(s,['Restricted crop: RandomResizedCrop(scale=(0.8, 1.0))','Horizontal flip retained for simple orientation variation','Normalize images before training','VerticalFlip, Rotation pipeline, and ColorJitter were tested but not used'],0.55,1.25,5.0,1.7,11)
if Path('images/augmented_samples.png').exists(): s.shapes.add_picture('images/augmented_samples.png', Inches(6.5), Inches(1.1), width=Inches(4.5))
rect=s.shapes.add_shape(1, Inches(0.55), Inches(5.35), Inches(7.8), Inches(0.58)); rect.fill.solid(); rect.fill.fore_color.rgb=GREEN; rect.line.color.rgb=RGBColor(187,247,208)
tx(s,'Decision: do not add preprocessing unless the ablation study shows a real improvement.',0.75,5.53,7.4,0.22,10,True,DARK,PP_ALIGN.CENTER); footer(s,4)

# 5 architecture table
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Custom CNN Architecture','Layer-by-layer operations, output sizes, and trainable parameter counts.')
rows=[['Stage','Operation','Output Size','Params','Role'],['Input','RGB image','3 x 64 x 64','0','Image tensor'],['Block 1','Conv 3->32 + BN + ReLU + MaxPool','32 x 32 x 32','960','Low-level edges'],['Block 2','Conv 32->64 + BN + ReLU + MaxPool','64 x 16 x 16','18,624','Texture patterns'],['Block 3','Conv 64->128 + BN + ReLU + MaxPool','128 x 8 x 8','74,112','Land-use features'],['Flatten','Reshape features','8,192','0','Vectorize'],['FC1','Linear 8192->256 + ReLU + Dropout','256','2,097,408','Classifier body'],['FC2','Linear 256->10','10 classes','2,570','Class scores']]
table=s.shapes.add_table(len(rows),5,Inches(0.55),Inches(1.2),Inches(6.9),Inches(4.4)).table
for r,row in enumerate(rows):
    for c,val in enumerate(row):
        cell=table.cell(r,c); cell.text=val; cell.margin_left=Inches(0.03); cell.margin_right=Inches(0.03)
        cell.text_frame.paragraphs[0].font.size=Pt(6.5 if r else 7)
        if r==0: cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(226,232,240); cell.text_frame.paragraphs[0].font.bold=True
metric(s,8.1,1.2,2.3,0.8,'Total trainable parameters','2,193,674')
rows2=[['Part','Params'],['Conv + BN blocks','93,696'],['FC1 layer','2,097,408'],['FC2 layer','2,570']]
tab=s.shapes.add_table(len(rows2),2,Inches(8.1),Inches(2.5),Inches(2.3),Inches(1.5)).table
for r,row in enumerate(rows2):
    for c,val in enumerate(row):
        tab.cell(r,c).text=val; tab.cell(r,c).text_frame.paragraphs[0].font.size=Pt(7)
        if r==0: tab.cell(r,c).fill.solid(); tab.cell(r,c).fill.fore_color.rgb=RGBColor(226,232,240); tab.cell(r,c).text_frame.paragraphs[0].font.bold=True
tx(s,'Interpretation',8.1,4.35,2.5,0.22,8,True); tx(s,'Most parameters are in FC1 after flattening. The convolution blocks learn spatial features; the classifier maps them to 10 land-use classes.',8.1,4.62,2.75,0.7,7,False,GRAY); footer(s,5)

# 6 baseline problem
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Baseline Problem: Training Setup Mattered','The first simple CNN was not bad, but the setup underperformed.')
s.shapes.add_picture(str(pipeline), Inches(0.8), Inches(1.3), width=Inches(5.2))
metric(s,8.2,1.5,1.45,0.8,'Baseline CNN','85.60%')
metric(s,10.0,1.5,1.6,0.8,'Incremental CNN','92.07%',GREEN)
bullet(s,['CNN architecture stayed the same.','Crop scale was restricted to preserve global structure.','Tuning increased from 4 to 6 epochs.','Accuracy improved by +6.47 percentage points.'],8.1,3.0,4.0,1.3,10)
footer(s,6)

# 7 tuning
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Hyperparameter Tuning','Experiment 2 was selected using lowest validation loss.')
s.shapes.add_picture(str(tuning), Inches(0.7), Inches(1.35), width=Inches(5.2))
rows=[['Exp','LR','Batch','Dropout','Val Loss'],['1','0.0010','32','0.30','0.4675'],['2','0.0005','32','0.30','0.4457'],['3','0.0010','64','0.50','0.6175']]
table=s.shapes.add_table(4,5,Inches(6.55),Inches(1.45),Inches(4.6),Inches(1.7)).table
for r,row in enumerate(rows):
    for c,val in enumerate(row):
        cell=table.cell(r,c); cell.text=val; cell.text_frame.paragraphs[0].font.size=Pt(8)
        if r==0: cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(226,232,240); cell.text_frame.paragraphs[0].font.bold=True
bullet(s,['Validation loss was the selection metric.','Experiment 2 gave the lowest validation loss.','Final training used LR=0.0005, batch=32, dropout=0.3.'],6.65,3.75,4.8,1.0,10)
footer(s,7)

# 8 comparison
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Final Model Comparison','The final Custom CNN outperformed the frozen ResNet18 baseline.')
s.shapes.add_picture(str(compare), Inches(0.8), Inches(1.35), width=Inches(5.3))
metric(s,6.65,1.6,1.6,0.8,'Final CNN Accuracy','92.07%',GREEN); metric(s,8.65,1.6,1.6,0.8,'ResNet18 Accuracy','82.44%')
bullet(s,['ResNet18 was used as a frozen feature extractor.','Only the final classification layer was trained.','EuroSAT satellite images differ from ImageNet natural images.','Future work could unfreeze later ResNet blocks.'],6.65,3.1,4.3,1.4,10)
footer(s,8)

# 9 confusion
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Final Confusion Matrix','Errors mainly occur between visually similar land-use classes.')
if Path('images/confusion_matrix.png').exists(): s.shapes.add_picture('images/confusion_matrix.png', Inches(1.25), Inches(1.05), width=Inches(5.1))
bullet(s,['Most classes are predicted reliably.','Remaining errors are mostly between agricultural and vegetation classes.','Linear structures such as River and Highway can also be visually similar.','Per-class evaluation is more informative than accuracy alone.'],7.0,1.65,4.4,1.8,10)
tx(s,'Final test accuracy: 92.07%',7.0,5.4,3.0,0.3,13,True,DARK); footer(s,9)

# 10 ablation
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Ablation Study: More Augmentation Was Not Better','Extra preprocessing was tested one method at a time.')
s.shapes.add_picture(str(ablation), Inches(0.85), Inches(1.35), width=Inches(5.1))
rows=[['Setting','Acc.','Change'],['Final baseline','92.07%','0.00 pp'],['+ ColorJitter','91.48%','-0.59 pp'],['+ VerticalFlip','89.83%','-2.25 pp'],['+ Rotation pipeline','67.60%','-24.47 pp']]
t=s.shapes.add_table(5,3,Inches(6.45),Inches(1.35),Inches(4.6),Inches(2.3)).table
for r,row in enumerate(rows):
    for c,val in enumerate(row):
        cell=t.cell(r,c); cell.text=val; cell.text_frame.paragraphs[0].font.size=Pt(8)
        if r==0: cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(226,232,240); cell.text_frame.paragraphs[0].font.bold=True
tx(s,'Rotation result is interpreted cautiously because the tested pipeline also included resize-and-crop operations.',6.45,4.25,4.6,0.6,8,False,GRAY); footer(s,10)

# 11 conclusion
s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); title(s,'Conclusion and Future Work','A simple 3-layer CNN achieved strong EuroSAT performance.')
metric(s,1.05,1.55,1.45,0.8,'Accuracy','92.07%',GREEN); metric(s,2.95,1.55,1.45,0.8,'F1-score','92.09%'); metric(s,4.85,1.55,1.45,0.8,'Best method','Simple aug.')
bullet(s,['Training setup mattered more than adding architecture complexity.','Restricted crop preserved useful satellite-image structure.','Extra augmentation was tested but not adopted.','Future work: fine-tune ResNet18, test larger image size, add Grad-CAM, and repeat with multiple seeds.'],1.05,3.0,7.0,1.5,11)
tx(s,'Final takeaway: more preprocessing was not automatically better; the best model used the simplest validated pipeline.',1.05,6.25,9.5,0.35,11,True,DARK,PP_ALIGN.CENTER)
footer(s,11)

prs.save(OUT)
print('Saved', OUT)
